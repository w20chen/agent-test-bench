#!/usr/bin/env python3
"""Build an editable PPTX deck for representative trace cases."""

from __future__ import annotations

import argparse
import json
import math
import textwrap
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "ink": "111827",
    "muted": "6B7280",
    "soft": "E5E7EB",
    "paper": "F8FAFC",
    "llm": "0EA5E9",
    "tool": "F97316",
    "search": "7C3AED",
    "fetch": "14B8A6",
    "exec": "F59E0B",
    "edit": "EF4444",
    "read": "22C55E",
    "test": "2563EB",
    "success": "16A34A",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--summary",
        type=Path,
        default=Path("reports/trace_case_ppt/scratch/case_summary.json"),
        help="Input summary JSON produced by scripts/summarize_trace_cases.py.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("reports/trace_case_ppt/output/trace_case_studies.pptx"),
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--preview-dir",
        type=Path,
        default=Path("reports/trace_case_ppt/scratch/previews"),
        help="Directory for lightweight PNG previews.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    summary = json.loads(args.summary.read_text(encoding="utf-8"))
    cases = summary["cases"]

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_cover(prs, cases)
    add_overview(prs, cases, summary["timing_note"])
    add_case_atlas(prs, cases)
    for case in cases:
        add_case_slide(prs, case)
    add_pattern_slide(prs, cases)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    render_lightweight_previews(summary, args.preview_dir)
    print(f"Wrote {args.output}")
    print(f"Wrote previews to {args.preview_dir}")


def add_cover(prs: Presentation, cases: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_text(
        slide,
        "Benchmark traces as compact stories",
        0.75,
        0.55,
        7.6,
        1.1,
        size=34,
        bold=True,
    )
    add_text(
        slide,
        "LLM turns, tool calls, and elapsed active time for eight representative cases",
        0.78,
        1.55,
        7.4,
        0.55,
        size=15,
        color=COLORS["muted"],
    )
    add_text(
        slide,
        "Each timeline uses recorded trace spans. No benchmark-specific tuning, no inferred timing.",
        0.78,
        6.55,
        8.6,
        0.35,
        size=9.5,
        color=COLORS["muted"],
    )

    x0, y0 = 0.85, 2.65
    width = 9.6
    for idx, case in enumerate(cases[:6]):
        y = y0 + idx * 0.43
        total = max(float(case["llm_s"]) + float(case["tool_s"]), 1.0)
        llm_w = width * float(case["llm_s"]) / total
        tool_w = width - llm_w
        add_bar(slide, x0, y, llm_w, 0.11, COLORS["llm"])
        add_bar(slide, x0 + llm_w, y, tool_w, 0.11, COLORS["tool"])
        add_text(slide, case["label"], x0 + width + 0.18, y - 0.06, 2.1, 0.25, size=8.5)

    add_legend(slide, 0.85, 5.55)
    add_big_ratio(slide, cases)


def add_overview(
    prs: Presentation, cases: list[dict[str, Any]], timing_note: str
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Overview")
    add_text(
        slide,
        "The cases separate into three trace shapes",
        0.65,
        0.55,
        8.2,
        0.55,
        size=25,
        bold=True,
    )
    add_text(
        slide,
        "Code repair mixes LLM and tools; Terminal-Bench is mostly long LLM reasoning; DeepResearch is tool-heavy web evidence gathering.",
        0.66,
        1.08,
        10.8,
        0.45,
        size=12.5,
        color=COLORS["muted"],
    )

    left, top = 0.8, 1.85
    bar_w = 7.3
    for idx, case in enumerate(cases):
        y = top + idx * 0.48
        total = max(float(case["llm_s"]) + float(case["tool_s"]), 1.0)
        llm_frac = float(case["llm_s"]) / total
        tool_frac = 1.0 - llm_frac
        add_text(slide, case["label"], left, y - 0.05, 2.65, 0.25, size=8.8)
        add_bar(slide, left + 2.9, y, bar_w * llm_frac, 0.15, COLORS["llm"])
        add_bar(
            slide,
            left + 2.9 + bar_w * llm_frac,
            y,
            bar_w * tool_frac,
            0.15,
            COLORS["tool"],
        )
        add_text(
            slide,
            f"{case['llm_s']:.0f}s LLM / {case['tool_s']:.0f}s tools",
            left + 2.9 + bar_w + 0.18,
            y - 0.045,
            1.35,
            0.22,
            size=7.8,
            color=COLORS["muted"],
        )

    add_legend(slide, 3.7, 6.06)
    add_text(slide, timing_note, 0.8, 6.52, 11.6, 0.34, size=8.5, color=COLORS["muted"])


def add_case_atlas(prs: Presentation, cases: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "All cases")
    add_text(
        slide,
        "Eight trace stories on one page",
        0.65,
        0.5,
        7.0,
        0.45,
        size=24,
        bold=True,
    )
    add_text(
        slide,
        "Each row shows active-time LLM/tool structure, recorded timing, key tools, and the case-level action.",
        0.66,
        0.98,
        10.8,
        0.3,
        size=10.5,
        color=COLORS["muted"],
    )

    add_text(slide, "case", 0.72, 1.45, 1.75, 0.2, size=7.5, bold=True, color=COLORS["muted"])
    add_text(slide, "mini trace", 2.55, 1.45, 2.5, 0.2, size=7.5, bold=True, color=COLORS["muted"])
    add_text(slide, "timing", 5.72, 1.45, 1.0, 0.2, size=7.5, bold=True, color=COLORS["muted"])
    add_text(slide, "important tools", 6.85, 1.45, 1.8, 0.2, size=7.5, bold=True, color=COLORS["muted"])
    add_text(slide, "what it is doing", 8.82, 1.45, 3.4, 0.2, size=7.5, bold=True, color=COLORS["muted"])

    for idx, case in enumerate(cases):
        y = 1.78 + idx * 0.58
        add_rule(slide, 0.72, y - 0.1, 11.85, COLORS["soft"], weight=0.5)
        add_text(slide, case["label"], 0.72, y, 1.6, 0.22, size=7.4, bold=True)
        add_text(slide, case["benchmark"], 0.72, y + 0.21, 1.6, 0.18, size=6.3, color=COLORS["muted"])
        draw_mini_trace(slide, case, 2.55, y + 0.08, 2.85)
        add_text(
            slide,
            f"L {case['llm_s']:.0f}s / T {case['tool_s']:.0f}s",
            5.68,
            y + 0.08,
            1.0,
            0.18,
            size=6.6,
            color=COLORS["muted"],
        )
        tool_text = ", ".join(short_tool(t["name"]) for t in case["top_tools"][:4])
        add_text(slide, tool_text, 6.85, y + 0.03, 1.78, 0.3, size=6.4, color=COLORS["ink"])
        add_text(slide, atlas_narrative(case), 8.82, y - 0.01, 3.65, 0.36, size=6.4, color=COLORS["muted"])

    add_legend(slide, 0.72, 6.62)


def add_case_slide(prs: Presentation, case: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, case["benchmark"])
    title = case["label"]
    add_text(slide, title, 0.62, 0.48, 7.8, 0.5, size=24, bold=True)
    status = "passed" if case.get("success") else "not passed"
    add_text(
        slide,
        f"{case['actions']} actions | {case['n_llm_calls']} LLM calls | "
        f"{case['n_tool_calls']} tool calls | {status}",
        0.64,
        0.96,
        8.7,
        0.3,
        size=10.5,
        color=COLORS["muted"],
    )

    add_metric(slide, "LLM", f"{case['llm_s']:.0f}s", 9.55, 0.52, COLORS["llm"])
    add_metric(slide, "Tools", f"{case['tool_s']:.0f}s", 10.75, 0.52, COLORS["tool"])
    add_metric(slide, "Elapsed", f"{case['elapsed_s']:.0f}s", 11.95, 0.52, COLORS["ink"])

    add_text(slide, case["narrative"], 0.72, 1.54, 11.7, 0.42, size=12.2)
    draw_timeline(slide, case, 0.82, 2.35, 11.7, 1.72)
    draw_highlights(slide, case, 0.82, 4.58)
    draw_top_tools(slide, case, 8.55, 4.58)
    add_text(
        slide,
        f"Timing source: {case['llm_timing_source']}. Timeline is gap-compressed active time.",
        0.82,
        6.82,
        9.8,
        0.22,
        size=7.8,
        color=COLORS["muted"],
    )


def add_pattern_slide(prs: Presentation, cases: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_kicker(slide, "Takeaway")
    add_text(slide, "Three compact patterns are enough for the talk", 0.65, 0.55, 9.5, 0.55, size=25, bold=True)
    patterns = [
        (
            "Code repair",
            "read -> probe -> edit -> test",
            "SWE-bench and SWE-rebench traces show mixed LLM/tool time with file inspection, code execution, edits, and validation.",
            COLORS["exec"],
        ),
        (
            "Terminal task",
            "reasoning-heavy command workflow",
            "Terminal-Bench examples can spend most active time in LLM calls while tools act as sparse environment probes.",
            COLORS["llm"],
        ),
        (
            "Deep research",
            "search -> fetch -> synthesize",
            "DeepResearch cases are web-tool dominated: searches and fetches become the visual spine of the trace.",
            COLORS["search"],
        ),
    ]
    for idx, (name, phrase, body, color) in enumerate(patterns):
        x = 0.85 + idx * 4.1
        add_rule(slide, x, 1.75, 2.0, color)
        add_text(slide, name, x, 1.98, 3.5, 0.34, size=17, bold=True)
        add_text(slide, phrase, x, 2.42, 3.4, 0.3, size=11, color=color, bold=True)
        add_text(slide, body, x, 2.92, 3.35, 1.1, size=10.5, color=COLORS["muted"])
    add_text(
        slide,
        "For slides, the trace does not need a full HTML replay. A two-lane active-time strip plus 3-5 labeled tools carries the case story.",
        0.85,
        5.5,
        10.7,
        0.5,
        size=15,
        bold=True,
    )


def draw_timeline(
    slide: Any, case: dict[str, Any], x: float, y: float, w: float, h: float
) -> None:
    lane_h = 0.22
    axis_y = y + 1.28
    add_text(slide, "LLM lane", x, y - 0.05, 1.0, 0.25, size=8.5, color=COLORS["muted"])
    add_text(slide, "Tool lane", x, y + 0.59, 1.0, 0.25, size=8.5, color=COLORS["muted"])
    add_rule(slide, x + 0.95, y + 0.1, w - 0.95, COLORS["soft"], weight=1)
    add_rule(slide, x + 0.95, y + 0.74, w - 0.95, COLORS["soft"], weight=1)

    end = max(float(case["timeline_end_s"]), 1.0)
    body_x = x + 0.95
    body_w = w - 0.95
    for span in case["timeline"]:
        start = float(span["start_s"]) / end
        duration = max(float(span["duration_s"]) / end, 0.0015)
        bx = body_x + body_w * start
        bw = max(body_w * duration, 0.015)
        if span["kind"] == "llm":
            add_bar(slide, bx, y, bw, lane_h, COLORS["llm"])
        else:
            add_bar(slide, bx, y + 0.64, bw, lane_h, tool_color(str(span["tool"])))

    for frac, label in ((0, "0s"), (0.5, f"{end/2:.0f}s"), (1, f"{end:.0f}s active")):
        tx = body_x + body_w * frac
        add_rule(slide, tx, axis_y - 0.05, 0.001, COLORS["soft"], weight=1)
        add_text(slide, label, tx - 0.22, axis_y + 0.03, 0.55, 0.22, size=7.3, color=COLORS["muted"])

    for item in case["highlights"][:4]:
        frac = float(item["start_s"]) / end
        hx = body_x + body_w * frac
        add_rule(slide, hx, y + 0.92, 0.001, COLORS["ink"], weight=1.2)


def draw_mini_trace(slide: Any, case: dict[str, Any], x: float, y: float, w: float) -> None:
    end = max(float(case["timeline_end_s"]), 1.0)
    lane_h = 0.055
    add_rule(slide, x, y + 0.02, w, COLORS["soft"], weight=0.5)
    add_rule(slide, x, y + 0.18, w, COLORS["soft"], weight=0.5)
    for span in case["timeline"]:
        start = float(span["start_s"]) / end
        duration = max(float(span["duration_s"]) / end, 0.002)
        bx = x + w * start
        bw = max(w * duration, 0.01)
        if span["kind"] == "llm":
            add_bar(slide, bx, y - 0.01, bw, lane_h, COLORS["llm"])
        else:
            add_bar(slide, bx, y + 0.15, bw, lane_h, tool_color(str(span["tool"])))


def draw_highlights(slide: Any, case: dict[str, Any], x: float, y: float) -> None:
    add_text(slide, "Important tool calls", x, y, 3.0, 0.25, size=11.5, bold=True)
    for idx, item in enumerate(case["highlights"][:5]):
        yy = y + 0.42 + idx * 0.34
        color = tool_color(str(item["tool"]))
        add_bar(slide, x, yy + 0.045, 0.16, 0.12, color)
        detail = clean_detail(str(item.get("detail") or ""))
        line = f"{item['label']}  {format_seconds(float(item['duration_s']))}"
        add_text(slide, line, x + 0.24, yy, 1.4, 0.22, size=8.8, bold=True)
        add_text(slide, detail, x + 1.52, yy, 5.55, 0.22, size=7.7, color=COLORS["muted"])


def draw_top_tools(slide: Any, case: dict[str, Any], x: float, y: float) -> None:
    add_text(slide, "Tool mix", x, y, 2.2, 0.25, size=11.5, bold=True)
    max_count = max((tool["count"] for tool in case["top_tools"]), default=1)
    for idx, tool in enumerate(case["top_tools"][:6]):
        yy = y + 0.42 + idx * 0.31
        label = short_tool(str(tool["name"]))
        add_text(slide, label, x, yy - 0.02, 0.78, 0.2, size=7.5)
        add_bar(slide, x + 0.88, yy, 1.25 * tool["count"] / max_count, 0.09, tool_color(str(tool["name"])))
        add_text(slide, str(tool["count"]), x + 2.2, yy - 0.035, 0.35, 0.18, size=7.2, color=COLORS["muted"])


def add_metric(slide: Any, label: str, value: str, x: float, y: float, color: str) -> None:
    add_text(slide, value, x, y, 1.0, 0.28, size=16, bold=True, color=color, align="center")
    add_text(slide, label, x, y + 0.31, 1.0, 0.2, size=7.5, color=COLORS["muted"], align="center")


def add_big_ratio(slide: Any, cases: list[dict[str, Any]]) -> None:
    total_llm = sum(float(case["llm_s"]) for case in cases)
    total_tool = sum(float(case["tool_s"]) for case in cases)
    ratio = total_tool / max(total_llm + total_tool, 1.0)
    add_text(slide, f"{ratio * 100:.0f}%", 9.35, 2.18, 2.6, 0.75, size=38, bold=True, color=COLORS["tool"], align="center")
    add_text(slide, "of active time is tools across these cases", 9.25, 2.92, 2.9, 0.48, size=11, color=COLORS["muted"], align="center")


def add_legend(slide: Any, x: float, y: float) -> None:
    add_bar(slide, x, y + 0.05, 0.26, 0.09, COLORS["llm"])
    add_text(slide, "LLM", x + 0.34, y - 0.02, 0.5, 0.18, size=8.5, color=COLORS["muted"])
    add_bar(slide, x + 0.95, y + 0.05, 0.26, 0.09, COLORS["tool"])
    add_text(slide, "Tools", x + 1.29, y - 0.02, 0.7, 0.18, size=8.5, color=COLORS["muted"])


def add_background(slide: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(COLORS["paper"])


def add_kicker(slide: Any, label: str) -> None:
    add_text(slide, label.upper(), 0.65, 0.24, 3.0, 0.2, size=8.5, bold=True, color=COLORS["muted"])


def add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    bold: bool = False,
    color: str = COLORS["ink"],
    align: str = "left",
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bar(slide: Any, x: float, y: float, w: float, h: float, color: str) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(max(w, 0.001)), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_rule(slide: Any, x: float, y: float, w: float, color: str, *, weight: float = 2) -> None:
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(weight)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def tool_color(tool: str) -> str:
    if tool in {"web_search"}:
        return COLORS["search"]
    if tool in {"web_fetch"}:
        return COLORS["fetch"]
    if tool in {"edit_file", "write_file"}:
        return COLORS["edit"]
    if tool in {"read_file", "list_dir"}:
        return COLORS["read"]
    if tool.startswith("exec-pytest"):
        return COLORS["test"]
    if tool.startswith("exec-"):
        return COLORS["exec"]
    return COLORS["tool"]


def short_tool(tool: str) -> str:
    return tool.replace("exec-", "").replace("_file", "").replace("web_", "")


def atlas_narrative(case: dict[str, Any]) -> str:
    text = str(case["narrative"])
    return textwrap.shorten(text, width=92, placeholder="...")


def clean_detail(detail: str) -> str:
    detail = " ".join(detail.replace("\\n", " ").split())
    detail = detail.replace('{"command": "', "").replace('", "working_dir": "/testbed"}', "")
    return textwrap.shorten(detail, width=76, placeholder="...")


def format_seconds(value: float) -> str:
    if 0.0 < value < 0.1:
        return "<0.1s"
    return f"{value:.1f}s"


def render_lightweight_previews(summary: dict[str, Any], preview_dir: Path) -> None:
    preview_dir.mkdir(parents=True, exist_ok=True)
    cases = summary["cases"]
    render_preview_cover(preview_dir / "slide_01_cover.png", cases)
    render_preview_overview(preview_dir / "slide_02_overview.png", cases)
    render_preview_atlas(preview_dir / "slide_03_all_cases.png", cases)
    for idx, case in enumerate(cases, start=4):
        render_preview_case(preview_dir / f"slide_{idx:02d}_{safe_name(case['label'])}.png", case)
    render_preview_takeaway(preview_dir / "slide_12_takeaway.png")


def new_preview() -> tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGB", (1920, 1080), "#" + COLORS["paper"])
    return img, ImageDraw.Draw(img)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        r"C:\Windows\Fonts\aptos.ttf",
        r"C:\Windows\Fonts\Aptos.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def render_preview_cover(path: Path, cases: list[dict[str, Any]]) -> None:
    img, draw = new_preview()
    draw.text((110, 95), "Benchmark traces as compact stories", fill="#" + COLORS["ink"], font=font(56, True))
    draw.text((112, 205), "LLM turns, tool calls, and elapsed active time for eight representative cases", fill="#" + COLORS["muted"], font=font(24))
    for idx, case in enumerate(cases[:6]):
        y = 390 + idx * 62
        total = max(float(case["llm_s"]) + float(case["tool_s"]), 1)
        llm_w = int(1080 * float(case["llm_s"]) / total)
        draw.rectangle((120, y, 120 + llm_w, y + 16), fill="#" + COLORS["llm"])
        draw.rectangle((120 + llm_w, y, 120 + 1080, y + 16), fill="#" + COLORS["tool"])
        draw.text((1230, y - 8), case["label"], fill="#" + COLORS["ink"], font=font(18))
    img.save(path)


def render_preview_overview(path: Path, cases: list[dict[str, Any]]) -> None:
    img, draw = new_preview()
    draw.text((94, 92), "The cases separate into three trace shapes", fill="#" + COLORS["ink"], font=font(42, True))
    for idx, case in enumerate(cases):
        y = 275 + idx * 66
        total = max(float(case["llm_s"]) + float(case["tool_s"]), 1)
        llm_w = int(1160 * float(case["llm_s"]) / total)
        draw.text((110, y - 12), case["label"], fill="#" + COLORS["ink"], font=font(18))
        draw.rectangle((520, y, 520 + llm_w, y + 22), fill="#" + COLORS["llm"])
        draw.rectangle((520 + llm_w, y, 1680, y + 22), fill="#" + COLORS["tool"])
        draw.text((1705, y - 8), f"{case['llm_s']:.0f}s / {case['tool_s']:.0f}s", fill="#" + COLORS["muted"], font=font(15))
    img.save(path)


def render_preview_atlas(path: Path, cases: list[dict[str, Any]]) -> None:
    img, draw = new_preview()
    draw.text((94, 82), "Eight trace stories on one page", fill="#" + COLORS["ink"], font=font(42, True))
    draw.text((96, 145), "Mini trace + timing + important tools + what the case is doing", fill="#" + COLORS["muted"], font=font(19))
    for idx, case in enumerate(cases):
        y = 235 + idx * 86
        draw.line((110, y - 15, 1810, y - 15), fill="#" + COLORS["soft"], width=1)
        draw.text((112, y), case["label"], fill="#" + COLORS["ink"], font=font(16, True))
        draw.text((112, y + 24), case["benchmark"], fill="#" + COLORS["muted"], font=font(12))
        draw_preview_mini_trace(draw, case, 470, y + 12, 420)
        draw.text((925, y + 12), f"L {case['llm_s']:.0f}s / T {case['tool_s']:.0f}s", fill="#" + COLORS["muted"], font=font(13))
        tools = ", ".join(short_tool(t["name"]) for t in case["top_tools"][:4])
        draw.text((1095, y + 12), tools, fill="#" + COLORS["ink"], font=font(13))
        draw.text((1365, y + 12), atlas_narrative(case), fill="#" + COLORS["muted"], font=font(12))
    img.save(path)


def draw_preview_mini_trace(draw: ImageDraw.ImageDraw, case: dict[str, Any], x: int, y: int, w: int) -> None:
    end = max(float(case["timeline_end_s"]), 1.0)
    draw.line((x, y, x + w, y), fill="#" + COLORS["soft"], width=1)
    draw.line((x, y + 26, x + w, y + 26), fill="#" + COLORS["soft"], width=1)
    for span in case["timeline"]:
        bx = x + int(w * float(span["start_s"]) / end)
        bw = max(int(w * float(span["duration_s"]) / end), 1)
        if span["kind"] == "llm":
            draw.rectangle((bx, y - 4, bx + bw, y + 4), fill="#" + COLORS["llm"])
        else:
            draw.rectangle((bx, y + 22, bx + bw, y + 30), fill="#" + tool_color(str(span["tool"])))


def render_preview_case(path: Path, case: dict[str, Any]) -> None:
    img, draw = new_preview()
    draw.text((95, 70), case["benchmark"].upper(), fill="#" + COLORS["muted"], font=font(15, True))
    draw.text((95, 112), case["label"], fill="#" + COLORS["ink"], font=font(40, True))
    draw.text((95, 195), case["narrative"], fill="#" + COLORS["ink"], font=font(22))
    draw_preview_timeline(draw, case, 130, 345, 1640)
    draw.text((130, 665), "Important tool calls", fill="#" + COLORS["ink"], font=font(22, True))
    for idx, item in enumerate(case["highlights"][:5]):
        y = 720 + idx * 42
        color = "#" + tool_color(str(item["tool"]))
        draw.rectangle((130, y + 8, 154, y + 24), fill=color)
        draw.text((170, y), f"{item['label']} {format_seconds(float(item['duration_s']))}", fill="#" + COLORS["ink"], font=font(17, True))
        draw.text((365, y), clean_detail(str(item.get("detail") or "")), fill="#" + COLORS["muted"], font=font(15))
    img.save(path)


def draw_preview_timeline(draw: ImageDraw.ImageDraw, case: dict[str, Any], x: int, y: int, w: int) -> None:
    draw.text((x, y - 38), "LLM lane", fill="#" + COLORS["muted"], font=font(15))
    draw.text((x, y + 74), "Tool lane", fill="#" + COLORS["muted"], font=font(15))
    body_x = x + 150
    body_w = w - 150
    draw.line((body_x, y, body_x + body_w, y), fill="#" + COLORS["soft"], width=2)
    draw.line((body_x, y + 112, body_x + body_w, y + 112), fill="#" + COLORS["soft"], width=2)
    end = max(float(case["timeline_end_s"]), 1.0)
    for span in case["timeline"]:
        bx = body_x + int(body_w * float(span["start_s"]) / end)
        bw = max(int(body_w * float(span["duration_s"]) / end), 2)
        if span["kind"] == "llm":
            draw.rectangle((bx, y - 15, bx + bw, y + 15), fill="#" + COLORS["llm"])
        else:
            draw.rectangle((bx, y + 95, bx + bw, y + 125), fill="#" + tool_color(str(span["tool"])))


def render_preview_takeaway(path: Path) -> None:
    img, draw = new_preview()
    draw.text((95, 105), "Three compact patterns are enough for the talk", fill="#" + COLORS["ink"], font=font(44, True))
    labels = ["Code repair", "Terminal task", "Deep research"]
    subtitles = ["read -> probe -> edit -> test", "reasoning-heavy command workflow", "search -> fetch -> synthesize"]
    colors = [COLORS["exec"], COLORS["llm"], COLORS["search"]]
    for idx, label in enumerate(labels):
        x = 130 + idx * 590
        draw.line((x, 300, x + 280, 300), fill="#" + colors[idx], width=8)
        draw.text((x, 340), label, fill="#" + COLORS["ink"], font=font(30, True))
        draw.text((x, 405), subtitles[idx], fill="#" + colors[idx], font=font(21, True))
    img.save(path)


def safe_name(value: str) -> str:
    return "".join(ch if ch.isalnum() or ch in "-_" else "_" for ch in value)[:60]


if __name__ == "__main__":
    main()
